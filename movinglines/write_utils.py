
# coding: utf-8

# In[1]:

def __init__():
    """
    Name:  

        write_utils

    Purpose:  

        Module regrouping codes that are used to write out certain tyes of files. Example: ict
        See each function within this module
        Contains the following functions:
            - write_ict: for writing out an ict file
            - merge_dicts: for merging multiple dicts
            - ict_tester: for testing the ict files

    Dependencies:

        - numpy
        - datetime

    Needed Files:

      None

    Modification History:

        Written: Samuel LeBlanc, NASA Ames, Santa Cruz, CA, 2016-03-25
    """
    pass


# In[366]:

def write_ict(header_dict,data_dict,filepath,data_id,loc_id,date,rev,order=[],default_format='.3f',file_comment=''):
    """
    Purpose:
        to write out a file in the ICARTT file format used by NASA archiving
    
    Input:
        filepath: full path of folder of the file to be saved
        data_id: for the prefix of the file name the instrumnet identifier
        loc_id: for the prefix of the file name, the location identifier
        date: date of the data, used in file naming convention
        rev: revision value of the file (RA,RB for infield, R0,R1,... for archiving)
        data_dict: dictionary with each key representing a different variable to be saved
                   each key is a dictionary of its own with the following keys:
                       - data: time series numpy array of the data to be saved
                       - unit: string value of the unit, if 'None' or '' set to be 'unitless' by default
                       - long_description: the long description of the variable
                       - format: (optional) the format used for writing out the data, ex: 4.2f
        header_dict: dictionary with a set of predefined keys for writing the header information, each is a string
                    - PI: name of PI
                    - Institution: name of the institution
                    - Instrument: full name of instrument/data source
                    - campaign: name of mission/campaign
                    - volume: (default 1) volume number
                    - file_num: (default 1) number of files
                    - time_interval: (default 1 second) the number of seconds between each archive
                    - indep_var_name: (defaults to Start_UTC) name of the data_dict variable which is the independant variable
                    - missing_val: (defaults to -9999) value that replaces the missing data
                    - special_comments: (optional) string of special comments to be written out. 
                                        Newline must be indicated by '\n'
                    - PI_contact: contact info for the PI
                    - platform: full platform name
                    - location: full name of location : refer to campaign location or plane location
                    - associated_data: (defaults to 'N/A') list files of associated data
                    - instrument_info: more detailed information on the instrument
                    - data_info: any specific info on the data, like time averging, ppm by volume
                    - uncertainty: specific notes about data uncertainty
                    - ULOD_flag: (defaults to -7777) value written when variable is past the upper limit of detection
                    - ULOD_value: (defaults to 'N/A') value of the upper limit of detection
                    - LLOD_flag: (defaults to -8888) value written when the measurement is below the lower limit of detection
                    - LLOD_value: (defaults to 'N/A') value of the lower limit of detection
                    - DM_contact: Contact information of the data manager for this data
                                  Name, affiliation, phone number, mailing address, email address and/or fax number.
                    - project_info: Information on the project
                    - stipulations: Details the stipulations on the use of this data
                    - Comments: (optional, if empty, returns 'N/A') any specific comments that is wanted to be included. 
                                if multiple lines, seperate with '\n'
                    - rev_comments: Comments related to each revision, newest first. Each ine is seperate by a '\n'
   
    Output:
        ict file with name format: '{data_id}_{loc_id}_{date}_{rev}.ict'
    
    Keywords: 
        order:(optional) list of names of the data variables put in order that will be saved in the file.
              if omitted, variables will be saved in random order.
        default_format: (defaults to '.3f') The format to use when writing out numbers 
                        when no specific format is defined for each data variable
        file_comment: (optional) If you want to put in a comment in the file name
        
    Dependencies:
        Numpy
        datetime
        write_utils (this module)
        
    Needed Files:
        None
        
    Example:
        see code for ict_tester function in this module
        
    Modification History:
        Written: Samuel LeBlanc, NASA Ames, Santa Cruz, 2016-03-25, Holy Friday
    """
    # module loads
    import numpy as np
    from datetime import datetime
    #from write_utils import merge_dicts
    # Should do input checking...
    
    # set up file path to write to
    data_id = data_id.replace('_','-')
    loc_id = loc_id.replace('_','-')
    if file_comment:
        file_comment = '_'+file_comment.replace(' ','-').replace('_','-')
    fname = filepath+'{data_id}_{loc_id}_{date}_{rev}{file_comment}.ict'.format(data_id=data_id,loc_id=loc_id,date=date,rev=rev,file_comment=file_comment)
    #f = open(fname,'w')
    
    # set the default dict values
    def_dict = {'rev':rev,
                'volume':1,
                'file_num':1,
                'time_interval':1.0,
                'indep_var_name':'Start_UTC',
                'missing_val':-9999,
                'special_comments':'',
                'associated_data':'N/A',
                'ULOD_flag':-7777,'ULOD_value':'N/A',
                'LLOD_flag':-8888,'LLOD_value':'N/A',
                'Comments':'N/A',
                'nlines':'{nlines}',
                'date_y':date[:4],'date_m':date[4:6],'date_d':date[6:],
                'now':datetime.now(),
                'num_data':len(data_dict)-1, # remove one for the independent variable
                'num_special_comments':len(header_dict['special_comments'].splitlines())}
    head = merge_dicts(def_dict,header_dict)
    
    # Compile the header information and verify some inputs
    head['data_head'] = ','.join(('1 '*head['num_data']).split())+'\n'+                               ','.join(('{missing_val} '*head['num_data']).split()).format(**head)
    head['data_format'] = '{t:.0f}'
    head['data_names'] = '{indep_var_name}'.format(**head)
    nv = head['indep_var_name']
    head['indep_var_unit'],head['indep_var_desc'] = data_dict[nv]['unit'],data_dict[nv]['long_description']
    head['rev_comments'] = head['rev_comments'].strip()
    if head['rev_comments'].find(head['rev'])<0:
        print("*** Revision comments don't include the current revision, please update ***")
        print('*** exiting, file not saved ***')
        return 
    if head['rev_comments'].find(head['rev'])>0:
        print("""*** Revision comments are not in the right order please update
    Have the current revision identifier in the top place ***""")
        print('*** exiting, file not saved ***')
        return 
    dnames = []
    if not order:
        order = data_dict.keys()            
    for n in order:
        print(n)
        if not n==head['indep_var_name']:
            stemp = '{n}, {unit}, {long_description}'.format(n=n,**data_dict[n])
            head['data_head'] = head['data_head']+'\n'+stemp
            if 'format' in data_dict[n]:
                fmt = data_dict[n]['format']
            else:
                fmt = default_format
            head['data_format'] = head['data_format']+',{:'+'{fmt}'.format(fmt=fmt)+'}'
            head['data_names'] = head['data_names']+','+n
            dnames.append(str(n))
    try:
        head['support_info'] = """-----------------------------------------------------------------------------
PI_CONTACT_INFO: {PI_contact}
PLATFORM: {platform}
LOCATION: {location}
ASSOCIATED_DATA: {associated_data}
INSTRUMENT_INFO: {instrument_info}
DATA_INFO: {data_info}
UNCERTAINTY: {uncertainty}
ULOD_FLAG: {ULOD_flag}
ULOD_VALUE: {ULOD_value}
LLOD_FLAG: {LLOD_flag}
LLOD_VALUE: {LLOD_value}
DM_CONTACT_INFO: {DM_contact}
PROJECT_INFO: {project_info}
STIPULATIONS_ON_USE: {stipulations}
OTHER_COMMENTS: {Comments}
REVISION: {rev}
{rev_comments}
-----------------------------------------------------------------------------
{data_names}""".format(**head)
    except KeyError as v:
        print('*** problem with header value of {v} ***'.format(v=v))
        print('*** exiting, file not saved ***')
        return
    head['num_info'] = len(head['support_info'].splitlines())
    try:
        head_str = """{nlines}, 1001
{PI}
{Institution}
{Instrument}
{campaign}
{volume},{file_num}
{date_y},{date_m},{date_d},{now:%Y,%m,%d}
{time_interval}
{indep_var_name}, {indep_var_unit}, {indep_var_desc}
{num_data}
{data_head}
{num_special_comments}
{special_comments}
{num_info}
{support_info}
""".format(**head)
    except KeyError as v:
        print('*** problem with header value of {v} ***'.format(v=v))
        print('*** exiting, file not saved ***')
        return
    
    # Now open and write out the header and data to the file
    with open(fname,'w') as f:
        f.write(head_str.format(nlines=len(head_str.splitlines())))
        for i,t in enumerate(data_dict[head['indep_var_name']]['data']):
            dat = [] # build each line and run checks on the data
            for n in dnames:
                d = data_dict[n]['data'][i]
                if not np.isfinite(d):
                    d = head['missing_val']
                if not type(head['ULOD_value']) is str:
                    if d>head['ULOD_value']:
                        d = head['ULOD_flag']
                if not type(head['LLOD_value']) is str:
                    if d<head['LLOD_value']:
                        d = head['LLOD_flag']
                dat.append(float(d))
            try:
                f.write(head['data_format'].format(*dat,t=t)+'\n')
            except:
                import pdb; pdb.set_trace()
    print('File writing successful to: {}'.format(fname))
    return


# In[265]:

def merge_dicts(*dict_args):
    """
    Given any number of dicts, shallow copy and merge into a new dict,
    precedence goes to key value pairs in latter dicts.
    """
    result = {}
    for dictionary in dict_args:
        result.update(dictionary)
    return result


# In[282]:

def ict_tester():
    """
    Simple function to test the write_ict file function
    makes a file with dummy variables
    """
    import numpy as np
    d_dict = {'Start_UTC':{'data':[230,231,232],'unit':'seconds from midnight UTC','long_description':'time keeping'},
          'X1':{'data':[1,2,3],'unit':'None','long_description':'test 1'},
          'X2':{'data':[10.9,11.9,12.9],'unit':'None','long_description':'test2'},
          'X3':{'data':[-2,-3,np.NaN],'unit':'somethinf','long_description':'tutor3'}
          }
    print(d_dict)
    hdict = {'PI':'Samuel LeBlanc',
         'Institution':'NASA Ames',
         'Instrument':'tester',
         'campaign':'NAAMES tester',
         'special_comments':'Only for testing with 3 data points',
         'PI_contact':'Samuel LeBlanc, samuel.leblanc@nasa.gov',
         'platform':'C130',
         'location':'based out of Santa Cruz, actual location in C130 file',
         'instrument_info':'None',
         'data_info':'made up',
         'uncertainty':'Undefined',
         'DM_contact':'See PI',
         'project_info':'NAAMES tester, made up data',
         'stipulations':'None',
         'rev_comments':"""  RA: first test of it\nR0: older"""
        }
    print(hdict)
    order = ['X1','X2','X3']
    write_ict(hdict,d_dict,filepath='C:/Users/sleblan2/Research/NAAMES/',
              data_id='4STAR_test',loc_id='C130',date='20160402',rev='RA',order=order)    


# In[369]:

def prep_data_for_ict(data_dict,Start_UTC=None,End_UTC=None,
                      in_var_name='utc',out_var_name='Start_UTC', in_input=True,time_interval=1.0):
    """
    Purpose:
        To create the time variable that matches the requirement of a ict file (ICARTT) for NASA archiving
        converts the in_var_name variable from utc hours, or datetime object to secdonds from utc save to the out_var_name
        Takes in the data_dict and makes sure each variable has data corresponding to each time entry
        Does nearest neighbor interpolation for linking each second measurement to a array value
        Creates an uninterupted data stream from the first point to the last, with missing data identifiers in between.
        Returns a modified data array.
    
    Input:
        data_dict: dictionary with each key representing a different variable to be saved
                   each key is a dictionary of its own with the following keys:
                       - data: time series numpy array of the data to be saved
                       - unit (not used)
                       - long_description (not used)
                       - format (not used)
                       
    Output:
        modified data array in data_dict, to have a continuous time series without time gaps
    
    Keywords: 
        in_var_name: (defaults to utc) the name of the variable holding the time series at the native resolution
        out_var_name: (defaults to Start_UTC) the name of the variable holding the seconds from midnight values
        Start_UTC: (optional) the start point of the time series, if not the first time point of the measurements
        End_UTC: (optional) the end point of the time series, if not the last point of the measurements
        in_input: (defaults to True) if set to True the Start_UTC and End_UTC use the same time series
                            units as the in_var_name, if False, uses the out_var_name (seconds from midnight)
        time_interval: (defaults to 1 second) the time interval of the nearest neighbor interpolation and to be saved
        
    Dependencies:
        Numpy
        datetime
        write_utils (this module)
        
    Needed Files:
        None
        
    Example:
        ...
        
    Modification History:
        Written: Samuel LeBlanc, NASA Ames, Santa Cruz, 2016-04-04
    """
    import numpy as np
    from write_utils import nearest_neighbor
    
    # check input
    if not in_var_name in data_dict:
        print("*** the variable defined by '{}' should be included in the data_dict ***".format(in_var_name))
    iv = in_var_name
    ov = out_var_name
    if type(data_dict[iv]['data']) is np.ndarray:
        if data_dict[iv]['data'].dtype is np.dtype(object):
            print('input variable not a recognized type')
            return
        elif data_dict[iv]['data'].dtype is np.dtype(float):
            # manageable type of float utc hours
            utcs = data_dict[iv]['data']*3600.0
    else:
        if type(data_dict[iv]['data']) is float:
            utcs = np.array(data_dict[iv]['data'])*3600.0
        else:
            print('non manageable input type, please make utc hours')
            return   
    
    # get the limits of the time series
    if not Start_UTC:
        Start_UTC = utcs[0]
    else:
        if in_input:
            Start_UTC = Start_UTC*3600.0
        
    if not End_UTC:
        End_UTC = utcs[-1]
    else:
        if in_input:
            End_UTC = End_UTC*3600.0
    
    # create the out_var_name array
    utc_out = np.arange(Start_UTC,End_UTC+1,time_interval)
    data_out = data_dict.copy()
    
    # now run through each data_dict variable to get the nearest neighbor
    for n in data_out:
        new = nearest_neighbor(utcs,data_out[n]['data'],utc_out,dist=time_interval/2.0)
        data_out[n]['data'] = new
    data_out[ov] = {'data':utc_out,'unit':'Seconds',
                    'long_description':'Time of measurement continuous starting from midnight UTC'}
    del(data_out[iv])
    return data_out


# In[330]:

def nearest_neighbor(X,Y,Xnew,dist=1):
    """
    Purpose:
        To return a nearest neighbor linear interpolation, but with a limit on the possible distance between the two points
    
    Input:
        X: initial independent variable
        Y: initial dependant variable
        Xnew: new independant variable to interpolate over
        dist: max distance allowed
                       
    Output:
        Ynew: new dependent variable interpolate using nearest neighbor
    
    Keywords: 
        dist: (default 1) see above
        
    Dependencies:
        Numpy
        Sp_parameters
        
    Example:
        ...
        
    Modification History:
        Written: Samuel LeBlanc, NASA Ames, Santa Cruz, 2016-04-04
    """
    from Sp_parameters import find_closest
    import numpy as np
    i = find_closest(X,Xnew)
    Ynew = Y[i]
    i_bad = abs(X[i]-Xnew) > dist
    try:
        Ynew[i_bad] = np.nan
    except ValueError:
        YYnew = Ynew.astype('float64')
        YYnew[i_bad] = np.nan
        Ynew = YYnew
    return Ynew   


# In[365]:

def make_plots_ict(data_dict,filepath,data_id,loc_id,date,rev,plot_together=[],plot_together2=[],indep_var_name='Start_UTC'):
    """
    Purpose:
        To plot the variables saved in the data_dict for prepping when saving
    
    Input:
        filepath: full path of folder of the file to be saved
        data_id: for the prefix of the file name the instrumnet identifier
        loc_id: for the prefix of the file name, the location identifier
        date: date of the data, used in file naming convention
        rev: revision value
        data_dict : see description in write_ict, uses data and unit
                       
    Output:
        plots
    
    Keywords:
        plot_together: (optional), list of names of variables to be plotted together on the same figure
        plot_together2: (optional), second list of names of variables to be plotted together on the same figure
        indep_var_name: (defaults to Start_UTC) the variable name of the independent variable in data_dict
        
    Dependencies:
        Numpy
        Sp_parameters
        matplotlib
        
    Example:
        ...
        
    Modification History:
        Written: Samuel LeBlanc, NASA Ames, Santa Cruz, 2016-04-04
    """
    import numpy as np
    import matplotlib.pyplot as plt
    plt.rc('text', usetex=False)
    utc = data_dict[indep_var_name]['data']
    ll = data_dict.keys()
    ll.remove(indep_var_name)
    if plot_together:
        fig = plt.figure()
        for n in plot_together:
            ll.remove(n)
            plt.plot(utc,data_dict[n]['data'],'x',label=n)
        plt.legend(frameon=False)
        plt.xlabel('UTC [seconds from midnight]')
        plt.ylabel('Values')
        plt.title(u'{data_id}_{loc_id}_{date}_{rev}.ict'.format(data_id=data_id,loc_id=loc_id,date=date,rev=rev))
        print('plotting the togethers')
        fig.savefig(filepath+'{data_id}_{loc_id}_{date}_{rev}_together.png'.format(                data_id=data_id,loc_id=loc_id,date=date,rev=rev),dpi=600,transparent=True)
        
    if plot_together2:
        fig = plt.figure()
        for n in plot_together2:
            ll.remove(n)
            plt.plot(utc,data_dict[n]['data'],'+',label=n)
        plt.legend(frameon=False)
        plt.xlabel('UTC [seconds from midnight]')
        plt.ylabel('Values')
        plt.title(u'{data_id}_{loc_id}_{date}_{rev}.ict'.format(data_id=data_id,loc_id=loc_id,date=date,rev=rev))
        print('plotting the togethers 2')
        fig.savefig(filepath+'{data_id}_{loc_id}_{date}_{rev}_together2.png'.format(                data_id=data_id,loc_id=loc_id,date=date,rev=rev),dpi=600,transparent=True)
        
    for n in ll:
        fig = plt.figure()
        plt.plot(utc,data_dict[n]['data'],'s',label=n)
        plt.legend(frameon=False)
        plt.xlabel('UTC [seconds from midnight]')
        plt.ylabel('{n} [{unit}]'.format(n=n,unit=data_dict[n].get('unit')))
        plt.title(u'{data_id}_{loc_id}_{date}_{rev} for {n}'.format(data_id=data_id,loc_id=loc_id,date=date,rev=rev,n=n))
        print('plotting {}'.format(n))
        fig.savefig(filepath+'{data_id}_{loc_id}_{date}_{rev}_{n}.png'.format(                data_id=data_id,loc_id=loc_id,date=date,rev=rev,n=n),dpi=600,transparent=True)

def create_generic_pptx(slides,filepath,title="My Presentation",subtitle=None):
    """
    Creates a PowerPoint presentation (.pptx) with custom slides.

    Args:
        slides (list of dict): Each dictionary represents a slide.
            - For text content: {"text": "Your slide content"}
            - For image content: {"image_path": "path/to/image.jpg"}
            - For bulelts content: {"bulelts":["bulelt1","bullet 2"]}
            - For multiple images: {"multiple_images":["path/to/image1.jpg"]}
            - For title of slide : {"title":"Your slide title"}
            
        title (str): Title for the presentation (default: "My Presentation")
        subtitle (str): subtitle for the presentation title slide
    """
    import os
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    prs = Presentation()

    # Add a title slide
    slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(slide_layout)
    title_placeholder = title_slide.placeholders[0]
    title_placeholder.text = title
    subtitle_placeholder = title_slide.placeholders[1]
    subtitle_placeholder.text = str(subtitle)
    
    left = Inches(0.5)
    top = Inches(1)
            
    # Add custom slides
    for slide_data in slides:
        slide_layout = prs.slide_layouts[6]  # Blank slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        if 'text' in slide_data and isinstance(slide_data['text'], list):
            # Concatenate text items with newline
            slide_data['text'] = '\n'.join(slide_data['text'])
        
        if "text" in slide_data:
            # Add text content
            textbox = slide.shapes.add_textbox(left, top, width=Inches(6), height=Inches(6))
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            p.text = slide_data["text"]
        if "image_path" in slide_data:
            # Add image content
            img_path = slide_data["image_path"]
            #left = (prs.slide_width - Inches(3)) / 2
            #top = (prs.slide_height - Inches(3)) / 2
            if 'text' in slide_data or 'bullets' in slide_data or 'table' in slide_data:
                pic = slide.shapes.add_picture(slide_data['image_path'], Inches(4.5), Inches(1.5), width=Inches(5.5))
                pic.top = int((prs.slide_height - pic.height) / 2)
            else:
                pic = slide.shapes.add_picture(img_path, 0, 0, width=Inches(10))
                pic.top = int((prs.slide_height - pic.height) / 2)
        if "bullets" in slide_data:
            # Add bullet points
            textbox = slide.shapes.add_textbox(left, top, width=Inches(6), height=Inches(6))
            text_frame = textbox.text_frame
            for bullet in slide_data["bullets"]:
                p = text_frame.add_paragraph()
                p.text = '- '+bullet
                p.level = 0
                p.space_after = Inches(0.03)  # Adjust spacing between bullets
        if 'table' in slide_data:
            rows = len(slide_data['table'])
            cols = len(slide_data['table'][0])

            table = slide.shapes.add_table(rows, cols, Inches(0.2), Inches(2), Inches(4.5), Inches(4)).table
            for i in range(rows):
                for j in range(cols):
                    table.cell(i, j).text = str(slide_data['table'][i][j])

                    # Apply formatting to the table cells (optional)
                    cell = table.cell(i, j)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    cell.text_frame.paragraphs[0].space_after = Inches(0.01)
                    if i==0: cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
               
        if 'multiple_images' in slide_data:
            num_images = len(slide_data['multiple_images'])
            if num_images ==2:
                i = 0
                for img_path in slide_data['multiple_images']:
                    pic = slide.shapes.add_picture(img_path, Inches(5)*i, top, width=Inches(5))
                    i += 1
            else:
                image_width = Inches(3)
                total_width = num_images * image_width
                left = (prs.slide_width - total_width) / 2
                top = Inches(1.5)
                for img_path in slide_data['multiple_images']:
                    slide.shapes.add_picture(img_path, left, top, width=image_width)
                    left += image_width
        
        if "title" in slide_data:
            # Add subtitle
            textbox = slide.shapes.add_textbox(Inches(2), Inches(0.2), width=Inches(8), height=Inches(1))
            text_frame = textbox.text_frame
            p = text_frame.add_paragraph()
            p.text = slide_data["title"]
            p.font.size = Pt(24)  # Customize font size for subtitles
            
    # Save the presentation
    prs.save(filepath)
    print("Presentation saved under the name: {}.".format(filepath))
